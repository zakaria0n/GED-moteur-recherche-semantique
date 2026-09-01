"""Extraction de texte : PDF, PDF scanne (OCR), Word, PowerPoint, Excel, images et GIF."""

import hashlib
import os
import re
import threading
from concurrent.futures import ThreadPoolExecutor, as_completed
from pathlib import Path

import numpy as np
import pypdfium2 as pdfium
import pypdfium2.raw as pdfium_c
from docx import Document
from openpyxl import load_workbook
from PIL import Image, ImageSequence
from pptx import Presentation
from pypdf import PdfReader


# Densite de rendu des pages PDF avant OCR (200 dpi : bon compromis vitesse/qualite).
OCR_RENDER_SCALE = 200 / 72

OCR_INSTANCE = None


def normalize_whitespace(text):
    text = text.replace("\r\n", "\n").replace("\r", "\n")
    text = re.sub(r"[ \t\f\v]+", " ", text)
    text = re.sub(r" *\n *", "\n", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def looks_garbled(text, min_letter_ratio=0.20):
    """Detecte une extraction PDF illisible (couche de texte presente mais
    corrompue) afin de declencher l'OCR a la place.

    On ne s'appuie que sur le ratio de lettres : un vrai texte court (titre,
    "AO 2025") reste a 100% de lettres et n'est jamais reclasse comme
    illisible, alors qu'une extraction brouillee (symboles, espaces) chute
    bien en dessous du seuil.
    """

    if not text:
        return False

    letters = sum(1 for char in text if char.isalpha())
    ratio = letters / max(1, len(text))

    return ratio < min_letter_ratio


def build_ocr():
    """Moteur OCR unique (RapidOCR/ONNX) : rapide sur CPU et stable."""
    global OCR_INSTANCE

    if OCR_INSTANCE is not None:
        return OCR_INSTANCE

    from rapidocr_onnxruntime import RapidOCR

    try:
        OCR_INSTANCE = RapidOCR()
        return OCR_INSTANCE
    except Exception as exc:
        raise ValueError("Impossible d'initialiser RapidOCR") from exc


def ocr_image_array(image):
    """Lance l'OCR sur une image numpy et renvoie le texte normalise."""
    results, _ = build_ocr()(image)
    if not results:
        return ""

    lines = [item[1] for item in results if item[1]]
    return normalize_whitespace("\n".join(lines))


def render_pdf_pages(pdf_path):
    """Rend chaque page du PDF en image numpy (pour l'OCR)."""
    doc = pdfium.PdfDocument(str(pdf_path))
    try:
        for page in doc:
            bitmap = page.render(scale=OCR_RENDER_SCALE)
            yield np.array(bitmap.to_pil().convert("RGB"))
    finally:
        doc.close()


def extract_text_from_pdf(pdf_path):
    path = Path(pdf_path)

    if not path.is_file():
        raise FileNotFoundError("PDF introuvable")

    try:
        reader = PdfReader(path)
    except Exception as exc:
        raise ValueError("Impossible de lire le PDF") from exc

    pages_text = []

    for page in reader.pages:
        text = page.extract_text()

        if text:
            pages_text.append(text)

    full_text = "\n\n".join(pages_text)

    return normalize_whitespace(full_text)


def extract_text_with_ocr(pdf_path):
    """OCR de chaque page d'un PDF scanne (rendu image puis RapidOCR)."""
    path = Path(pdf_path)

    if not path.is_file():
        raise FileNotFoundError("PDF introuvable")

    try:
        pages_text = []
        for image in render_pdf_pages(path):
            text = ocr_image_array(image)
            if text:
                pages_text.append(text)
    except ValueError:
        raise
    except Exception as exc:
        raise ValueError("Impossible de lire le PDF scanne avec OCR") from exc

    return normalize_whitespace("\n\n".join(pages_text))


def ocr_pdf_pages(pdf_path):
    """OCR de chaque page rendue : capture le texte des images imbriquees.

    Renvoie le texte combine (normalise) de toutes les pages. Leve une
    ValueError si le rendu/OCR echoue, attrape par l'appelant.
    """

    pages_text = []

    for image in render_pdf_pages(pdf_path):
        text = ocr_image_array(image)
        if text:
            pages_text.append(text)

    return normalize_whitespace("\n\n".join(pages_text))


def _image_fingerprint(image_array):
    """Empreinte rapide d'une image pour la deduplication.

    Redimensionne a 32x32, convertit en niveaux de gris, et hash le
    resultat. Deux images visuellement identiques (meme logo) auront
    le meme hash meme si la compression differe legerement.
    """

    try:
        pil_image = Image.fromarray(image_array).convert("L").resize((32, 32))
        pixels = list(pil_image.getdata())
        return hashlib.md5(bytes(pixels)).hexdigest()
    except Exception:
        return None


def _iter_embedded_images(pdf_path, max_pages=None):
    """Genere (page_index, largeur_px, hauteur_px, PIL.Image) pour chaque image
    raster imbreee du PDF, via l'API objets de pypdfium2.

    Attention : PdfPage n'expose PAS de get_images() — la liste des images
    passe par get_objects(filter=[FPDF_PAGEOBJ_IMAGE]). Les images dont le
    bitmap est illisible (encodage exotique) sont ignorees.
    """

    try:
        document = pdfium.PdfDocument(str(pdf_path))
    except Exception:
        return

    try:
        page_count = len(document) if max_pages is None else min(len(document), max_pages)

        for page_index in range(page_count):
            try:
                objects = document[page_index].get_objects(filter=[pdfium_c.FPDF_PAGEOBJ_IMAGE])
                for obj in objects:
                    try:
                        try:
                            pil_image = obj.get_bitmap(render=False).to_pil()
                        except Exception:
                            pil_image = obj.get_bitmap(render=True).to_pil()

                        yield page_index, pil_image
                    except Exception:
                        continue
            except Exception:
                continue
    finally:
        document.close()


def _extract_page_images(pdf_path, page_index):
    """Extrait les images raster d'une page donnee (numpy RGB).

    Les petites images (< 1000 pixels de surface) sont ignorees
    (decorations, puces, separateurs).
    """

    result = []

    for img_page, pil_image in _iter_embedded_images(pdf_path):
        if img_page != page_index:
            continue

        width, height = pil_image.size

        if width * height < 1000:
            continue

        result.append(np.array(pil_image.convert("RGB")))

    return result


def _collect_unique_images(pdf_path):
    """Parcourt le PDF et collecte les images uniques (dedup par hash).

    Renvoie (pages_with_images, unique_images) :
      - pages_with_images : {page_index: [image_arrays]} (uniques)
      - unique_images : liste plate des images uniques a OCRiser
    """

    pages_with_images = {}
    seen_fingerprints = set()
    unique_images = []

    for page_index, pil_image in _iter_embedded_images(pdf_path):
        try:
            width, height = pil_image.size

            # Ignorer les petites images (decorations, icones).
            if width * height < 1000:
                continue

            image_array = np.array(pil_image.convert("RGB"))

            # Deduplication par hash perceptuel : le meme logo sur 50 pages
            # n'est conserve qu'une seule fois.
            fingerprint = _image_fingerprint(image_array)

            if fingerprint is None:
                continue

            if fingerprint in seen_fingerprints:
                continue

            seen_fingerprints.add(fingerprint)
            unique_images.append(image_array)
            pages_with_images.setdefault(page_index, []).append(image_array)
        except Exception:
            continue

    return pages_with_images, unique_images


def ocr_pdf_images_only(pdf_path):
    """OCR intelligent : deduplique les images avant OCR, en parallele.

    Au lieu d'OCRiser chaque page (ou chaque image), on detecte les
    images en double (logo d'entreprise, tampon, etc.) et on n'OCRise
    qu'une seule copie de chaque image unique. Pour un PDF de 50 pages
    avec le meme logo en en-tete, c'est 50x plus rapide.
    Le cache est GLOBAL au processus : une image deja vue dans un autre
    document du corpus n'est pas re-OCRisee (papiers en-tete partages).

    Renvoie le texte combine (normalise) des images uniques uniquement.
    """

    pages_with_images, unique_images = _collect_unique_images(pdf_path)

    if not unique_images:
        return ""

    return ocr_images_parallel(unique_images)


# --- OCR parallele ---

_DEFAULT_WORKERS = min(8, (os.cpu_count() or 4))

# --- Cache OCR global des images (multi-documents) ---
# Un meme logo/tampon present dans plusieurs PDFs (papier en-tete d'entreprise)
# n'est OCRise qu'une seule fois par session : la premiere rencontre remplit le
# cache (empreinte -> texte), les suivantes le relisent sans OCR.
_OCR_IMAGE_CACHE: dict[str, str] = {}
_OCR_IMAGE_CACHE_MAX = 1000
_OCR_IMAGE_CACHE_LOCK = threading.Lock()


def _ocr_image_cached(image_array):
    """OCR d'une image avec cache global par empreinte visuelle."""

    fingerprint = _image_fingerprint(image_array)

    if fingerprint is not None:
        with _OCR_IMAGE_CACHE_LOCK:
            cached = _OCR_IMAGE_CACHE.get(fingerprint)

        if cached is not None:
            return cached

    text = ocr_image_array(image_array)

    if fingerprint is not None and text:
        with _OCR_IMAGE_CACHE_LOCK:
            if len(_OCR_IMAGE_CACHE) >= _OCR_IMAGE_CACHE_MAX:
                _OCR_IMAGE_CACHE.pop(next(iter(_OCR_IMAGE_CACHE)))
            _OCR_IMAGE_CACHE[fingerprint] = text

    return text


def _ocr_single_page(image):
    """OCR d'une seule page/image (wrapper pour ThreadPoolExecutor)."""
    try:
        return _ocr_image_cached(image)
    except Exception:
        return ""


def ocr_pdf_pages_parallel(pdf_path, max_workers=None):
    """OCR parallele de chaque page rendue.

    Le rendu des pages (pypdfium2) est sequentiel (pas thread-safe),
    mais l'OCR RapidOCR/ONNX est parallélise via un pool de threads
    (ONNX libere le GIL).

    Renvoie le texte combine (normalise) de toutes les pages.
    """
    max_workers = max_workers or _DEFAULT_WORKERS

    # Phase 1 : rendu sequentiel (pypdfium2 n'est pas thread-safe).
    pages = list(render_pdf_pages(pdf_path))

    if not pages:
        return ""

    if len(pages) == 1:
        text = _ocr_single_page(pages[0])
        return normalize_whitespace(text) if text else ""

    # Phase 2 : OCR parallele des images rendues.
    pages_text = [None] * len(pages)

    with ThreadPoolExecutor(max_workers=min(max_workers, len(pages))) as pool:
        futures = {
            pool.submit(_ocr_single_page, img): i
            for i, img in enumerate(pages)
        }
        for future in as_completed(futures):
            idx = futures[future]
            try:
                pages_text[idx] = future.result()
            except Exception:
                pages_text[idx] = ""

    result = [t for t in pages_text if t]
    return normalize_whitespace("\n\n".join(result))


def ocr_images_parallel(images, max_workers=None):
    """OCR parallele d'une liste d'images numpy.

    Utile pour ocr_pdf_images_only : OCRise toutes les images uniques
    en parallele au lieu d'une par une.
    """
    max_workers = max_workers or _DEFAULT_WORKERS

    if not images:
        return ""

    if len(images) == 1:
        # Passer quand meme par le cache : une image unique ici peut avoir
        # deja ete OCRisee dans un autre document (logo partage).
        text = _ocr_single_page(images[0])
        return normalize_whitespace(text) if text else ""

    pages_text = []

    with ThreadPoolExecutor(max_workers=min(max_workers, len(images))) as pool:
        futures = [pool.submit(_ocr_single_page, img) for img in images]
        for future in as_completed(futures):
            try:
                text = future.result()
                if text:
                    pages_text.append(text)
            except Exception:
                continue

    return normalize_whitespace("\n\n".join(pages_text))


def has_raster_images(pdf_path, max_pages=5):
    """Detecte rapidement si un PDF contient des images raster significatives.

    Sert de filtre ultra-rapide pour decider si l'OCR des images est
    necessaire. Un PDF vectoriel pur (texte + dessins vectoriels) n'a pas
    d'images raster et ne necessite PAS d'OCR.

    Renvoie True des qu'au moins une image raster significative (> 200x200 px)
    est detectee dans les premieres pages (max_pages).
    """

    try:
        for _, pil_image in _iter_embedded_images(pdf_path, max_pages=max_pages):
            width, height = pil_image.size

            if width * height > 40_000:
                return True

        return False
    except Exception:
        return True  # En cas de doute, on OCR quand meme


def count_embedded_images(pdf_path):
    """Compte les images raster imbriquees dans le PDF (meilleure-effort).

    Sert au monitoring et a la decision "OCR malin" : un PDF avec images mais
    peu de texte merite un passage OCR. En cas d'API indisponible, renvoie 0.
    """

    try:
        return sum(1 for _ in _iter_embedded_images(pdf_path))
    except Exception:
        return 0


def merge_text_layers(base, addition):
    """Fusionne deux extractions en conservant l'ordre de la couche de base.

    Ajoute uniquement les lignes de l'OCR absentes de la couche texte (comparaison
    normalisee insensible a la casse). Evite de dupliquer le texte deja present
    dans la couche texte tout en recuperant le texte des images.
    """

    if not addition:
        return base

    if not base:
        return addition

    base_lines = [line.strip() for line in base.splitlines() if line.strip()]
    base_lower = {line.lower() for line in base_lines}

    merged = list(base_lines)
    addition_lines = [line.strip() for line in addition.splitlines() if line.strip()]

    for line in addition_lines:
        if line.lower() not in base_lower:
            merged.append(line)
            base_lower.add(line.lower())

    return normalize_whitespace("\n".join(merged))


def extract_text_from_powerpoint(pptx_path):
    path = Path(pptx_path)

    if not path.is_file():
        raise FileNotFoundError("Fichier PowerPoint introuvable")

    try:
        presentation = Presentation(path)
    except Exception as exc:
        raise ValueError("Impossible de lire le fichier PowerPoint") from exc

    parts = []

    for slide_index, slide in enumerate(presentation.slides, start=1):
        slide_parts = []

        for shape in slide.shapes:
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    cell_values = []

                    for cell in row.cells:
                        text = cell.text.strip()

                        if text:
                            cell_values.append(text)

                    if cell_values:
                        slide_parts.append(" | ".join(cell_values))
            elif hasattr(shape, "text"):
                text = shape.text.strip()

                if text:
                    slide_parts.append(text)

        if slide_parts:
            parts.append(f"Slide {slide_index}")
            parts.append("\n".join(slide_parts))

    full_text = "\n\n".join(parts)

    return normalize_whitespace(full_text)


def extract_text_from_excel(xlsx_path):
    path = Path(xlsx_path)

    if not path.is_file():
        raise FileNotFoundError("Fichier Excel introuvable")

    try:
        workbook = load_workbook(path, data_only=True)
    except Exception as exc:
        raise ValueError("Impossible de lire le fichier Excel") from exc

    parts = []

    for sheet in workbook.worksheets:
        sheet_lines = []

        for row in sheet.iter_rows(values_only=True):
            values = []

            for cell in row:
                if cell is not None:
                    text = str(cell).strip()

                    if text:
                        values.append(text)

            if values:
                sheet_lines.append(" | ".join(values))

        if sheet_lines:
            parts.append(f"Sheet: {sheet.title}")
            parts.append("\n".join(sheet_lines))

    full_text = "\n\n".join(parts)

    return normalize_whitespace(full_text)


def extract_text_from_image(image_path):
    path = Path(image_path)

    if not path.is_file():
        raise FileNotFoundError("Image introuvable")

    try:
        with Image.open(path) as image:
            array = np.array(image.convert("RGB"))
    except Exception as exc:
        raise ValueError("Impossible de lire l'image") from exc

    try:
        return ocr_image_array(array)
    except Exception as exc:
        raise ValueError("Impossible de lire l'image avec OCR") from exc


def extract_text_from_gif(gif_path):
    path = Path(gif_path)

    if not path.is_file():
        raise FileNotFoundError("GIF introuvable")

    try:
        with Image.open(path) as image:
            frames = [frame.copy().convert("RGB") for frame in ImageSequence.Iterator(image)]
    except Exception as exc:
        raise ValueError("Impossible de lire le GIF") from exc

    if not frames:
        return ""

    # GIF long : on ne lit que 3 images representatives (debut, milieu, fin).
    if len(frames) <= 3:
        frame_indexes = list(range(len(frames)))
    else:
        frame_indexes = sorted({0, len(frames) // 2, len(frames) - 1})

    frame_texts = []
    seen_texts = set()

    for index in frame_indexes:
        try:
            text = ocr_image_array(np.array(frames[index]))
        except Exception as exc:
            raise ValueError("Impossible de lire le GIF avec OCR") from exc

        if text and text not in seen_texts:
            seen_texts.add(text)
            frame_texts.append(text)

    return normalize_whitespace("\n\n".join(frame_texts))


def extract_text_from_word(word_path):
    path = Path(word_path)

    if not path.is_file():
        raise FileNotFoundError("Fichier Word introuvable")

    try:
        document = Document(path)
    except Exception as exc:
        raise ValueError("Impossible de lire le fichier Word") from exc

    paragraphs = []

    for paragraph in document.paragraphs:
        text = paragraph.text.strip()

        if text:
            paragraphs.append(text)

    full_text = "\n\n".join(paragraphs)

    return normalize_whitespace(full_text)
